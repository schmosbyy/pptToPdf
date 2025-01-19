import zipfile
import os
import xml.etree.ElementTree as ET
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx import Presentation
from extract_video_url import extract_video_urls_from_pptx
import random
import re
from libreoffice import convert_to_pdf 

# Directory paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
LIBREOFFICE_DIR = os.path.join(BASE_DIR, "libreoffice")

# Helper function to convert EMUs to Inches
def emu_to_inches(emu):
    return emu / 914400.0

def list_media_shapes(slide):
    """
    List only media shapes (e.g., videos or multimedia elements) in a slide, 
    including their position, width, and height.
    """
    media_shapes_info = []

    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        # Check if the shape is of type MEDIA
        if shape.shape_type == 16:  # MEDIA type (e.g., video)
            media_shapes_info.append({
                'Type': shape.shape_type,
                'Position': (shape.left, shape.top),
                'Width': shape.width,
                'Height': shape.height
            })

    return media_shapes_info

def extract_video_position_from_slide(pptx_zip, slide_num):
    """
    Extract the position (left, top, width, height) of the video from the slide's XML.
    """
    slide_file = f"ppt/slides/slide{slide_num}.xml"
    try:
        with pptx_zip.open(slide_file) as slide_content:
            slide_tree = ET.parse(slide_content)
            slide_root = slide_tree.getroot()

            # Namespace for drawing elements
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # Iterate through all shapes in the slide
            for shape in slide_root.findall('.//a:sp', ns):
                
                # Check if the shape contains a video or multimedia element
                video_element = shape.find('.//a:video', ns)
                if video_element is not None:
                    
                    # Extract position and size of the shape (if available)
                    position = shape.find('.//a:spPr/a:xfrm', ns)
                    left = float(position.get('x', 0)) if position is not None else 0
                    top = float(position.get('y', 0)) if position is not None else 0
                    
                    # Get width and height of the shape (in EMUs)
                    width = float(shape.find('.//a:spPr/a:xfrm/a:ext', ns).get('cx', 0))
                    height = float(shape.find('.//a:spPr/a:xfrm/a:ext', ns).get('cy', 0))

                    # Convert EMUs to inches
                    left = emu_to_inches(left)
                    top = emu_to_inches(top)
                    width = emu_to_inches(width)
                    height = emu_to_inches(height)

                    return left, top, width, height
    except Exception as e:
        return None

def overlay_video_with_shape(pptx_directory, output_directory, extracted_video_urls):
    # Ensure the output directory exists
    output_directory = os.path.abspath(output_directory)
    if not os.path.exists(output_directory):
        os.makedirs(output_directory)

    # Process each PPTX file from the extracted video URLs
    for pptx_file, video_data in extracted_video_urls.items():
        pptx_path = os.path.join(pptx_directory, pptx_file)
        presentation = Presentation(pptx_path)
        output_file = os.path.join(output_directory, pptx_file)
        
        # Iterate through the slides and overlay a transparent rectangle where video is found
        for slide_index, video_url in video_data:
            
            # Extract the slide number from the string (if it's the XML path)
            match = re.search(r'slide(\d+)', slide_index)
            if match:
                slide_num = int(match.group(1))  # Extract and convert the slide number
            else:
                continue

            # Then use it to access the slide
            slide = presentation.slides[slide_num - 1]  # 1-based to 0-based index
                    
            # List and print all media shapes for the current slide
            media_shapes_info = list_media_shapes(slide)
            # print(f"Media shapes found on slide {slide_num}: {media_shapes_info}")
            
            # Iterate over each media shape and overlay the transparent rectangle using its position, width, and height
            for media in media_shapes_info:


                left, top = media['Position']
                width, height = media['Width'], media['Height']
                
                # Add the transparent rectangle using the position, width, and height of the media shape
                shape = slide.shapes.add_shape(
                    1,  # Rectangle shape
                    Inches(left / 914400), Inches(top / 914400), Inches(width / 914400), Inches(height / 914400)
                )
                shape.fill.background()  # Transparent fill
                
                # Add an isosceles triangle to represent the play button
                triangle_width = Inches(1)  # Width of the triangle (play button)
                triangle_height = Inches(1)  # Height of the triangle (play button)
                
                # Calculate the position to center the triangle inside the rectangle
                triangle_left = Inches(left / 914400) + Inches(width / 914400)/2 - triangle_width/2
                triangle_top = Inches(top / 914400) + Inches(height / 914400 )/2 - triangle_height/2
                
                # Create the isosceles triangle shape (play button)
                play_button = slide.shapes.add_shape(
                    7,  # Isosceles triangle shape
                    triangle_left, triangle_top, triangle_width, triangle_height
                )
                play_button.fill.solid()  # Solid color for the triangle
                play_button.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color for the play button (can change color)
                
                # Rotate the triangle 90 degrees to make it point to the right
                play_button.rotation = 90  # Rotate the triangle to the right

                # Set hyperlink to the extracted video URL for the current slide
                video_link = video_url
                link = shape.click_action
                link.hyperlink.address = video_link

        # Save the modified presentation
        presentation.save(output_file)

def main():
    pptx_directory = "ppt"
    output_directory = "output"

    extracted_video_urls = extract_video_urls_from_pptx(pptx_directory)

    overlay_video_with_shape(pptx_directory, output_directory, extracted_video_urls)

    # After processing the PPTX files, call the convert_to_pdf function from libreoffice_conversion.py
    for filename in os.listdir(output_directory):
        if filename.endswith(".pptx"):
            input_path = os.path.join(output_directory, filename)
            convert_to_pdf(input_path)  # Calling the function to convert the PPTX to PDF

    print(f"All conversions completed. PDFs saved in {output_directory}")

if __name__ == "__main__":
    main()